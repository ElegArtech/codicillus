"""L'EXTRACTION DES TROIS FORMATS BUREAUTIQUES — lot T-052.

Ce module est exécuté DANS UN SOUS-PROCESSUS, jamais dans le processus qui sert
les requêtes. C'est la raison d'être de tout le service, et elle est écrite noir
sur blanc par `STACK-TECHNIQUE.md` §4.6 : « ces convertisseurs sont lents,
consomment de la mémoire de façon irrégulière et échouent sur des fichiers
malformés. Les isoler garantit RG-M12-04 : un fichier en erreur ne peut pas
interrompre le lot, ni faire tomber l'application. »

Un fichier malformé peut donc ici tout se permettre — lever, boucler, saturer la
mémoire, tuer son interpréteur : il n'emporte que ce processus-là, et le service
répond quand même. `service.py` le lance avec un délai et récolte son verdict.

═══════════════════════════════════════════════════════════════════════════
CE QUE CE MODULE PRODUIT, ET CE QU'IL NE PRODUIRA JAMAIS — ADR-004

Il produit du **Markdown** et des **images extraites**. Rien d'autre.

`ADR-004` interdit nommément « toute conversion du document canonique effectuée
hors de l'application — en particulier dans le service Python, qui s'arrête à la
production de Markdown ». Aucun nom de nœud du format canonique n'apparaît dans
ce fichier, et c'est vérifiable en le lisant : la sortie est du texte, et c'est
l'application qui lui applique son convertisseur unique.

═══════════════════════════════════════════════════════════════════════════
UN SEUL ÉCRIVAIN DE MARKDOWN DANS LE SERVICE, ET C'EST PANDOC

`STACK` §4.6 attribue un outil par format : Pandoc lit le `.docx`, python-pptx
lit le `.pptx` — « Pandoc ne lit pas le pptx, vérifié dans son manuel, ce format
n'existe qu'en sortie » —, pdfplumber lit le `.pdf`. Ces trois outils tiennent
ici leur rôle exact : la LECTURE.

L'ÉCRITURE du Markdown, elle, est faite par Pandoc dans les trois cas. Les deux
extracteurs Python rendent du HTML — un balisage dont l'échappement est trivial
et sans ambiguïté (trois caractères) — et Pandoc le met en Markdown.

Le motif est le même que celui d'`ADR-004` un cran plus bas : écrire du Markdown
à la main demande de connaître les caractères qui y sont actifs, ceux qui ne le
sont qu'en tête de ligne, ceux qui ne le sont qu'en paire. Recopier ces règles
ici en ferait une seconde écriture, qui divergerait de la première au premier cas
limite — une astérisque dans un nom de serveur, un tiret bas dans un chemin. Un
seul écrivain, donc, et il est déjà dans l'image.

Conséquence assumée et déclarée au rapport du lot : Pandoc absent rend les trois
formats indisponibles, et non le seul `.docx`. Le contrôle de santé rend les
trois versions d'outils, l'application les lit, et un service amputé se voit.

═══════════════════════════════════════════════════════════════════════════
LE DIALECTE DE SORTIE EST CHOISI, PAS SUBI

`--to=gfm-raw_html --wrap=none --columns=1`, et chaque moitié compte :

  * `gfm` écrit des titres à dièses, des puces à tiret, des blocs clôturés par
    des accents graves et des tableaux à barres verticales avec leur ligne de
    filet — les formes que lit `src/lib/contenu/markdown.ts`, l'implémentation
    unique. Le dialecte `markdown` par défaut de Pandoc écrirait des titres
    soulignés, que cette implémentation ne lit pas comme des titres.
  * `-raw_html` : sans lui, une image légendée d'un `.docx` sort en bloc HTML de
    figure — cinq lignes de balises que l'implémentation unique lirait en
    paragraphes de texte, balises comprises (mesuré). Privé de HTML brut, Pandoc
    dégrade la figure en la forme d'image de Markdown, qui EST lue. La règle
    vaut plus largement : l'implémentation unique ne lit aucun HTML, donc rien
    de ce que le service produit ne doit en porter.
  * `--wrap=none` : l'implémentation unique lit UN PARAGRAPHE PAR LIGNE. Un
    repli de colonne à 72 caractères, qui est le défaut de Pandoc, ferait autant
    de paragraphes que de lignes repliées.
  * `--columns=1` : Pandoc aligne les colonnes d'un tableau en bourrant les
    cellules d'espaces, et le découpage de cellules de l'implémentation unique
    n'en retire QU'UNE de chaque bord. Sans cette option, chaque cellule d'un
    tableau importé porterait des espaces de fin — mesuré : `| Serveur    |`
    devient le texte « Serveur␣␣␣ ». La valeur 1 n'a pas d'effet de repli, que
    `--wrap=none` a déjà fermé.

CE QUE PANDOC NE SAIT PAS RENDRE, ET QUI EST DONC PERDU. Un bloc de code sans
langage est écrit par Pandoc en bloc INDENTÉ, dans tous ses dialectes Markdown
et sans option pour l'en empêcher (mesuré sur `gfm` et `commonmark`) ; or
l'implémentation unique ne lit comme bloc de code que les blocs clôturés. Le
texte survit — en paragraphes indentés —, sa nature de bloc de code non. Le
cahier des charges exige pour le `.docx` « titres, listes, tableaux, mise en
forme préservés » et ne nomme pas les blocs de code : la limite est constatée et
comptée au rapport du lot, elle n'est pas contournée par une réécriture du
Markdown de Pandoc.

═══════════════════════════════════════════════════════════════════════════
LES MOTIFS D'ÉCHEC SONT UN JEU FERMÉ, ET CHACUN A SA SOURCE

Aucun n'est inventé : les deux premiers sont les cas d'échec que le lot
d'exemple du gel de V-24 montre, avec leur formulation française
(`seeds/corpus.ts`, `LOT_IMPORT` — « le fichier est protégé par un mot de
passe », « le fichier semble endommagé : sa structure interne n'a pas pu être
ouverte »). Le service rend des CODES, jamais des phrases : la mise en français
appartient à l'interface, et STACK §4.7 lui réserve son catalogue.

Usage, depuis `service.py` :

    python -m convertisseurs <chemin> <format> <dossier-de-travail>

Le verdict part sur la sortie standard, en JSON, et le code de retour vaut 0 y
compris quand la conversion échoue : un échec de conversion est un RÉSULTAT,
pas une panne. Un code de retour non nul signale que ce processus est mort en
route, et `service.py` le traite comme un fichier endommagé.
"""

from __future__ import annotations

import base64
import html
import json
import mimetypes
import subprocess
import sys
from pathlib import Path
from typing import Any

# Les trois formats de la table de STACK §4.6 sont ceux de `CONVERTISSEURS`, plus
# bas, et il n'y en a pas de quatrième — la table d'aiguillage est la SEULE liste
# de ce module, plutôt qu'une liste doublée d'un tuple qui divergerait.
# Le `.md` et le `.txt` ne passent JAMAIS ici : « ne sort pas de l'application,
# c'est le chemin qui garantit l'idempotence et la résolution des références
# (RG-M12-01) ».

# L'avertissement d'un PDF sans texte extractible. LA PHRASE EST CELLE DU CAHIER
# DES CHARGES, M12.1, mot pour mot : « un PDF scanné produit une note avec un
# avertissement explicite "contenu scanné — transcription manuelle
# recommandée" ». Elle n'est pas rédigée ici, elle est recopiée — et la
# reconnaissance de caractères reste hors périmètre (STACK §4.6).
AVERTISSEMENT_SCANNE = "contenu scanné — transcription manuelle recommandée"

# Le code que porte cet avertissement dans la réponse, pour que l'application
# puisse le compter sans lire du français.
CODE_SCANNE = "contenu-scanne"

# Le dossier où les médias extraits sont déposés, et le préfixe sous lequel le
# Markdown les référence. Le nom rendu dans la réponse est le chemin RELATIF au
# dossier de travail — celui-là même que le Markdown référence, à ceci près que
# Pandoc préfixe le sien du dossier courant : `./media/…` référence
# `media/…`. L'appariement se fait donc sur le chemin normalisé, ce que fait
# toute résolution de chemin relatif (RG-M12-07, « les images référencées en
# chemin relatif sont reprises »).
DOSSIER_MEDIAS = "media"


class EchecDeConversion(Exception):
    """Un fichier que le service ne sait pas lire — un résultat, pas une panne."""

    def __init__(self, motif: str) -> None:
        super().__init__(motif)
        self.motif = motif


# ═══════════════════════════════════════════════════ Pandoc, l'écrivain ══


def _pandoc(arguments: list[str], entree: bytes | None, travail: Path) -> str:
    """Pandoc, lancé sur un fichier ou sur une entrée standard.

    AUCUN DÉLAI N'EST REPRIS ICI, et ce n'est pas un oubli : `service.py` borne
    et tue CE processus-ci, et Pandoc — son enfant — part avec lui. Un second
    délai imbriqué ferait deux bornes dont l'une mentirait sur l'autre.
    """
    try:
        acheve = subprocess.run(
            ["pandoc", *arguments],
            input=entree,
            capture_output=True,
            cwd=travail,
            check=False,
        )
    except FileNotFoundError as absent:
        raise EchecDeConversion("outil-absent") from absent
    if acheve.returncode != 0:
        raise EchecDeConversion("fichier-endommage")
    return acheve.stdout.decode("utf-8")


def _markdown_depuis_html(fragment: str, travail: Path) -> str:
    """Le Markdown d'un fragment HTML produit par l'un des deux extracteurs."""
    return _pandoc(
        ["--from=html", "--to=gfm-raw_html", "--wrap=none", "--columns=1"],
        fragment.encode("utf-8"),
        travail,
    )


def _medias(travail: Path) -> list[dict[str, Any]]:
    """Les fichiers déposés par l'extraction de médias, lus et encodés.

    Le nom rendu est le chemin RELATIF au dossier de travail — voir
    `DOSSIER_MEDIAS` pour ce qui l'apparie au chemin écrit dans le Markdown.
    """
    racine = travail / DOSSIER_MEDIAS
    if not racine.is_dir():
        return []
    sortis: list[dict[str, Any]] = []
    for chemin in sorted(p for p in racine.rglob("*") if p.is_file()):
        octets = chemin.read_bytes()
        nom = chemin.relative_to(travail).as_posix()
        sortis.append(
            {
                "nom": nom,
                "type_mime": mimetypes.guess_type(nom)[0] or "application/octet-stream",
                "octets": len(octets),
                "contenu_base64": base64.b64encode(octets).decode("ascii"),
            }
        )
    return sortis


# ═════════════════════════════════════════════════════════════ Le .docx ══


def convertir_docx(fichier: Path, travail: Path) -> dict[str, Any]:
    """`.docx` par Pandoc — STACK §4.6, « médias extraits ».

    `--extract-media` est relatif au dossier de travail : le Markdown rendu
    référence donc les images par un chemin relatif court, et non par le chemin
    absolu d'un dossier éphémère qui n'existera plus au retour.
    """
    markdown = _pandoc(
        [
            "--from=docx",
            "--to=gfm-raw_html",
            "--wrap=none",
            "--columns=1",
            "--extract-media=.",
            str(fichier),
        ],
        None,
        travail,
    )
    return {"markdown": markdown, "images": _medias(travail), "avertissements": []}


# ═════════════════════════════════════════════════════════════ Le .pptx ══


def _fragment_html(fragment: Any) -> str:
    """Un fragment de texte, sa graisse et son italique conservés.

    M12.1 attend du `.pptx` une « conversion en contenu riche » : le gras et
    l'italique d'un fragment sont la mise en forme que le format porte, et les
    perdre serait rendre du texte nu là où il y avait du texte formé.
    """
    texte = html.escape(fragment.text)
    if texte == "":
        return ""
    police = getattr(fragment, "font", None)
    if getattr(police, "italic", None):
        texte = f"<em>{texte}</em>"
    if getattr(police, "bold", None):
        texte = f"<strong>{texte}</strong>"
    return texte


def _elements_de_liste(cadre: Any) -> list[tuple[int, str]]:
    """Les paragraphes d'un cadre de texte : leur niveau, et leur HTML.

    Un cadre de texte de diapositive est une liste à puces : c'est ce qu'il est
    dans l'outil de présentation, et `paragraph.level` en porte l'imbrication.
    """
    sortis: list[tuple[int, str]] = []
    for paragraphe in cadre.paragraphs:
        contenu = "".join(_fragment_html(fragment) for fragment in paragraphe.runs)
        if contenu.strip() == "":
            continue
        sortis.append((min(int(getattr(paragraphe, "level", 0) or 0), 8), contenu))
    return sortis


def _liste_html(elements: list[tuple[int, str]]) -> list[str]:
    """UNE SEULE LISTE POUR TOUTE LA DIAPOSITIVE, ET C'EST DÉLIBÉRÉ.

    Deux listes HTML adjacentes font écrire à Pandoc un commentaire HTML de
    séparation entre les deux — l'implémentation unique le lirait comme un
    paragraphe portant ce commentaire en clair (mesuré). Les cadres de texte
    d'une diapositive sont donc réunis en une liste, ce qu'ils sont pour le
    lecteur de la diapositive.
    """
    if not elements:
        return []
    lignes: list[str] = []
    ouverts = 0
    for niveau, contenu in elements:
        vise = niveau + 1
        while ouverts < vise:
            lignes.append("<ul>")
            ouverts += 1
        while ouverts > vise:
            lignes.append("</ul>")
            ouverts -= 1
        lignes.append(f"<li>{contenu}</li>")
    lignes.extend(["</ul>"] * ouverts)
    return lignes


def _tableau_html(table: Any) -> list[str]:
    """Un tableau de diapositive, en HTML. La première ligne fait l'en-tête."""
    lignes = ["<table>"]
    for rang, ligne in enumerate(table.rows):
        balise = "th" if rang == 0 else "td"
        cellules = "".join(
            f"<{balise}>{html.escape(cellule.text.strip())}</{balise}>" for cellule in ligne.cells
        )
        lignes.append(f"<tr>{cellules}</tr>")
    lignes.append("</table>")
    return lignes


def convertir_pptx(fichier: Path, travail: Path) -> dict[str, Any]:
    """`.pptx` par python-pptx — « une section par diapositive » (M12.1).

    LA SECTION EST UN TITRE QUAND LA DIAPOSITIVE EN A UN, ET UN SÉPARATEUR
    SINON. Une diapositive sans titre n'en a pas, et lui en fabriquer un —
    « Diapositive 4 » — serait écrire un texte que ni le cahier des charges, ni
    la pile, ni le gel ne portent. Le séparateur marque la frontière de section
    sans rien inventer. Compté au rapport du lot comme un vide de la source.

    Le titre est de NIVEAU 2 : le niveau 1 est celui du titre de la note, que
    l'import tire du nom du fichier ou de l'en-tête (RG-M12-05).
    """
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError, PythonPptxError

    try:
        presentation = Presentation(str(fichier))
    except PackageNotFoundError as absente:
        raise EchecDeConversion("fichier-endommage") from absente
    except PythonPptxError as invalide:
        raise EchecDeConversion("fichier-endommage") from invalide
    except (KeyError, ValueError, OSError) as illisible:
        # Une archive ouvrable dont le contenu ne tient pas : le format
        # `.pptx` est un conteneur, et un conteneur peut être vide de sens.
        raise EchecDeConversion("fichier-endommage") from illisible

    dossier_medias = travail / DOSSIER_MEDIAS
    fragments: list[str] = []
    rang_image = 0

    for rang, diapositive in enumerate(presentation.slides, start=1):
        titre = None
        cadre_du_titre = None
        if diapositive.shapes.title is not None:
            cadre_du_titre = diapositive.shapes.title
            texte = cadre_du_titre.text.strip()
            titre = texte if texte != "" else None

        fragments.append(f"<h2>{html.escape(titre)}</h2>" if titre else "<hr/>")

        # Les cadres de texte de la diapositive sont accumulés, puis rendus en
        # UNE liste — voir `_liste_html`. Les tableaux et les images gardent
        # leur place dans l'ordre du document : ils interrompent la liste, ce
        # qui est fidèle, et n'en créent pas deux adjacentes.
        elements: list[tuple[int, str]] = []
        for forme in diapositive.shapes:
            if cadre_du_titre is not None and forme is cadre_du_titre:
                continue
            if getattr(forme, "has_table", False):
                fragments.extend(_liste_html(elements))
                elements = []
                fragments.extend(_tableau_html(forme.table))
                continue
            image = getattr(forme, "image", None)
            if image is not None:
                fragments.extend(_liste_html(elements))
                elements = []
                rang_image += 1
                dossier_medias.mkdir(parents=True, exist_ok=True)
                nom = f"diapositive-{rang}-{rang_image}.{image.ext}"
                (dossier_medias / nom).write_bytes(image.blob)
                alternative = html.escape((forme.name or "").strip())
                fragments.append(
                    f'<p><img src="{DOSSIER_MEDIAS}/{nom}" alt="{alternative}" /></p>'
                )
                continue
            if getattr(forme, "has_text_frame", False):
                elements.extend(_elements_de_liste(forme.text_frame))
        fragments.extend(_liste_html(elements))

    markdown = _markdown_depuis_html("\n".join(fragments), travail)
    return {"markdown": markdown, "images": _medias(travail), "avertissements": []}


# ══════════════════════════════════════════════════════════════ Le .pdf ══


def motif_pdf(erreur: BaseException) -> str:
    """LE MOTIF D'UN PDF QUI NE S'OUVRE PAS — protégé, ou endommagé.

    Les deux cas sont ceux que le lot d'exemple du gel de V-24 montre, et ils ne
    se disent pas la même chose à l'utilisateur : un fichier protégé se rouvre
    avec son mot de passe, un fichier endommagé ne se rouvre pas.

    LA DISTINCTION SE LIT DANS LA CAUSE, PAS DANS LE TYPE. pdfplumber enveloppe
    toute erreur de pdfminer dans une exception unique et pose l'originale en
    contexte : chercher `PDFEncryptionError` sur le type rendrait tout PDF
    chiffré « endommagé ». Mesuré sur un fichier au dictionnaire de chiffrement
    greffé.

    Cette fonction est isolée POUR AVOIR UN CAS D'ÉPREUVE (`P-26`) : elle se
    juge sur une exception fabriquée, sans PDF ni disque.
    """
    from pdfminer.pdfdocument import PDFEncryptionError

    vue: BaseException | None = erreur
    while vue is not None:
        if isinstance(vue, PDFEncryptionError):
            return "fichier-protege"
        vue = vue.__context__ if vue.__context__ is not vue else None
    return "fichier-endommage"


def convertir_pdf(fichier: Path, travail: Path) -> dict[str, Any]:
    """`.pdf` par pdfplumber — le texte SÉLECTIONNABLE, et rien de deviné.

    Aucun texte extractible : la note est créée quand même, avec l'avertissement
    du cahier des charges en tête. C'est le point le plus explicite de M12.1, et
    il ne se traduit ni par un échec, ni par une reconnaissance de caractères —
    « hors périmètre » (STACK §4.6).

    UNE LIGNE EXTRAITE FAIT UN PARAGRAPHE, ET C'EST UN CHOIX DÉCLARÉ. Un PDF ne
    porte pas de paragraphes : il porte des lignes posées à des coordonnées, et
    `extract_text` les rend séparées par un saut de ligne, sans ligne vide entre
    les paragraphes d'origine. Les recoller demanderait une heuristique
    d'interligne ou d'indentation qu'aucune source du dépôt ne décrit — et une
    heuristique fausse coupe des phrases en deux. Les lignes sont donc rendues
    telles qu'elles sont extraites. Compté au rapport du lot.
    """
    import pdfplumber

    try:
        with pdfplumber.open(str(fichier)) as document:
            pages = [page.extract_text() or "" for page in document.pages]
    except Exception as illisible:  # noqa: BLE001
        # LARGE PAR NÉCESSITÉ. pdfplumber enveloppe, pdfminer lève sa propre
        # famille, et un PDF malformé fait aussi lever les bibliothèques de
        # décodage d'image en dessous. Le motif, lui, reste précis.
        raise EchecDeConversion(motif_pdf(illisible)) from illisible

    avertissements: list[str] = []
    fragments: list[str] = []

    if all(page.strip() == "" for page in pages):
        avertissements.append(CODE_SCANNE)
        fragments.append(f"<p>{html.escape(AVERTISSEMENT_SCANNE)}</p>")

    for page in pages:
        for ligne in page.split("\n"):
            texte = ligne.strip()
            if texte != "":
                fragments.append(f"<p>{html.escape(texte)}</p>")

    markdown = _markdown_depuis_html("\n".join(fragments), travail)
    return {"markdown": markdown, "images": [], "avertissements": avertissements}


# ═════════════════════════════════════════════════════════════ L'aiguillage ══

CONVERTISSEURS = {
    "docx": convertir_docx,
    "pptx": convertir_pptx,
    "pdf": convertir_pdf,
}


def convertir(fichier: Path, format_: str, travail: Path) -> dict[str, Any]:
    """Le verdict d'un fichier : son Markdown, ses images, ses avertissements."""
    convertisseur = CONVERTISSEURS.get(format_)
    if convertisseur is None:
        raise EchecDeConversion("format-non-pris-en-charge")
    resultat = convertisseur(fichier, travail)
    return {"format": format_, **resultat}


def principal(arguments: list[str]) -> int:
    """Le point d'entrée du sous-processus. Il rend du JSON, toujours.

    AUCUNE TRACE TECHNIQUE NE SORT SUR LA SORTIE STANDARD — STACK §4.7 :
    « aucune trace technique ne remonte à l'interface ». Le détail part sur la
    sortie d'erreur, que `service.py` journalise et ne transmet pas.
    """
    if len(arguments) != 3:
        print(json.dumps({"issue": "echec", "motif": "format-non-pris-en-charge"}))
        return 0
    chemin, format_, travail = arguments
    try:
        resultat = convertir(Path(chemin), format_, Path(travail))
    except EchecDeConversion as echec:
        print(json.dumps({"issue": "echec", "motif": echec.motif}))
        return 0
    except MemoryError:
        # Un fichier qui demande plus que la machine n'a. Il est nommé à part
        # parce qu'il n'est pas endommagé : il est trop gros pour ce service.
        print(json.dumps({"issue": "echec", "motif": "fichier-trop-lourd"}), flush=True)
        return 0
    except BaseException as imprevu:  # noqa: BLE001
        # LE FILET DE RG-M12-04, ET IL EST DÉLIBÉRÉMENT LARGE. Un fichier
        # malformé qui fait lever une bibliothèque tierce d'une exception qu'elle
        # ne documente pas ne doit pas devenir une réponse vide : il devient un
        # fichier endommagé, consigné, et le lot continue.
        print(f"{type(imprevu).__name__}: {imprevu}", file=sys.stderr)
        print(json.dumps({"issue": "echec", "motif": "fichier-endommage"}))
        return 0
    print(json.dumps({"issue": "converti", **resultat}))
    return 0


if __name__ == "__main__":
    sys.exit(principal(sys.argv[1:]))
